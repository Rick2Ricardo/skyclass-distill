from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "output"
ASSETS = OUTPUT / "assets" / "figures"
PPTX_PATH = OUTPUT / "final_presentation_cn.pptx"

SW = 13.333
SH = 7.5

FONT_CN = "Arial Unicode MS"
FONT_EN = "Arial"

NAVY = "14213D"
INK = "1F2933"
SLATE = "5B6673"
MUTED = "7C8793"
PAPER = "FAF9F6"
WHITE = "FFFFFF"
SOFT = "F1F3F5"
LINE = "D7DCE1"
TEAL = "2A9D8F"
TEAL_LIGHT = "DCEFEA"
CORAL = "D86657"
CORAL_LIGHT = "F8E3DF"
GOLD = "D5A021"
GOLD_LIGHT = "F8EDCE"
BLUE = "3678B8"
BLUE_LIGHT = "DFEBF6"
PURPLE = "6B5CA5"
PURPLE_LIGHT = "E9E5F4"


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_font(font, name, size=None, color=None, bold=None, italic=None):
    font.name = name
    if size is not None:
        font.size = Pt(size)
    if color is not None:
        font.color.rgb = rgb(color)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic
    rpr = font._element
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = OxmlElement("a:ea")
        rpr.append(ea)
    ea.set("typeface", name)
    return font


def set_bg(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=16,
    color=INK,
    bold=False,
    font=FONT_CN,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    italic=False,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    set_font(p.font, font, size=size, color=color, bold=bold, italic=italic)
    p.line_spacing = line_spacing
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=INK, gap=6, bullet_color=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        set_font(p.font, FONT_CN, size=size, color=color)
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
    return box


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill=WHITE,
    line=LINE,
    radius=True,
    line_width=1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_picture_contain(slide, path, x, y, w, h):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw = iw * scale
    ph = ih * scale
    px = x + (w - pw) / 2
    py = y + (h - ph) / 2
    pic = slide.shapes.add_picture(
        str(path), Inches(px), Inches(py), Inches(pw), Inches(ph)
    )
    return pic


def add_source(slide, text, x=0.72, y=6.88, w=10.8, align=PP_ALIGN.LEFT):
    return add_text(
        slide,
        text,
        x,
        y,
        w,
        0.2,
        size=7.5,
        color=MUTED,
        font=FONT_EN if text.startswith("Source") else FONT_CN,
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, num):
    add_line(slide, 0.65, 7.14, 12.68, 7.14, color=LINE, width=0.7)
    add_text(
        slide,
        "Resource2Skill · 论文汇报",
        0.67,
        7.18,
        2.7,
        0.17,
        size=7.2,
        color=MUTED,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        f"{num:02d}",
        12.25,
        7.18,
        0.38,
        0.17,
        size=7.2,
        color=MUTED,
        font=FONT_EN,
        align=PP_ALIGN.RIGHT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_standard_title(slide, title, num, eyebrow=None):
    set_bg(slide)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.66),
        Inches(0.32),
        Inches(0.10),
        Inches(0.48),
    ).fill.solid()
    accent = slide.shapes[-1]
    accent.fill.fore_color.rgb = rgb(CORAL)
    accent.line.fill.background()
    if eyebrow:
        add_text(
            slide,
            eyebrow.upper(),
            0.92,
            0.20,
            3.0,
            0.18,
            size=7.5,
            color=CORAL,
            bold=True,
            font=FONT_EN,
            valign=MSO_ANCHOR.MIDDLE,
        )
    add_text(
        slide,
        title,
        0.91,
        0.40 if eyebrow else 0.31,
        11.7,
        0.58,
        size=25.5,
        color=NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, num)


def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text.strip()


def add_metric(slide, value, label, x, y, w, color=CORAL, value_size=28):
    add_text(
        slide,
        value,
        x,
        y,
        w,
        0.48,
        size=value_size,
        color=color,
        bold=True,
        font=FONT_EN,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        label,
        x,
        y + 0.45,
        w,
        0.42,
        size=10.5,
        color=SLATE,
        valign=MSO_ANCHOR.TOP,
    )


def add_chevron(slide, x, y, color=MUTED):
    add_text(
        slide,
        "→",
        x,
        y,
        0.35,
        0.35,
        size=20,
        color=color,
        font=FONT_EN,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def horizontal_bars(
    slide,
    data,
    x,
    y,
    w,
    h,
    min_value,
    max_value,
    highlight_index=0,
    label_w=1.95,
    value_w=0.75,
):
    row_h = h / len(data)
    bar_x = x + label_w
    bar_w = w - label_w - value_w
    for idx, (label, value) in enumerate(data):
        cy = y + idx * row_h
        color = TEAL if idx == highlight_index else BLUE
        light = TEAL_LIGHT if idx == highlight_index else BLUE_LIGHT
        add_text(
            slide,
            label,
            x,
            cy + 0.02,
            label_w - 0.10,
            row_h - 0.04,
            size=11.5,
            color=NAVY if idx == highlight_index else INK,
            bold=idx == highlight_index,
            font=FONT_EN,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_rect(
            slide,
            bar_x,
            cy + row_h * 0.25,
            bar_w,
            row_h * 0.50,
            fill=SOFT,
            line=SOFT,
            radius=True,
        )
        frac = max(0, min(1, (value - min_value) / (max_value - min_value)))
        add_rect(
            slide,
            bar_x,
            cy + row_h * 0.25,
            max(0.04, bar_w * frac),
            row_h * 0.50,
            fill=color,
            line=color,
            radius=True,
        )
        add_text(
            slide,
            f"{value:.1f}",
            x + w - value_w + 0.02,
            cy + 0.01,
            value_w - 0.02,
            row_h - 0.03,
            size=12,
            color=color if idx == highlight_index else SLATE,
            bold=True,
            font=FONT_EN,
            align=PP_ALIGN.RIGHT,
            valign=MSO_ANCHOR.MIDDLE,
        )


def build_deck():
    OUTPUT.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    # 1 — Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, NAVY)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.72),
        Inches(0.72),
        Inches(0.12),
        Inches(1.15),
    ).fill.solid()
    accent = slide.shapes[-1]
    accent.fill.fore_color.rgb = rgb(CORAL)
    accent.line.fill.background()
    add_text(
        slide,
        "RESOURCE2SKILL",
        1.06,
        0.64,
        3.7,
        0.32,
        size=10,
        color="93C5BD",
        bold=True,
        font=FONT_EN,
    )
    add_text(
        slide,
        "从人类创作资源中\n蒸馏智能体可执行技能",
        1.02,
        1.02,
        7.1,
        1.72,
        size=31,
        color=WHITE,
        bold=True,
        line_spacing=0.95,
    )
    add_text(
        slide,
        "Distilling Executable Skills from Human-Created Resources for Software Agents",
        1.06,
        2.88,
        6.8,
        0.72,
        size=12,
        color="C7CFD8",
        font=FONT_EN,
    )
    add_metric(slide, "+11.9 pp", "相对无技能智能体的平均提升", 9.15, 1.18, 3.1, color="F0B7AE", value_size=34)
    add_text(
        slide,
        "26 / 28",
        9.16,
        2.47,
        2.4,
        0.45,
        size=24,
        color="93C5BD",
        bold=True,
        font=FONT_EN,
    )
    add_text(
        slide,
        "模型×领域单元超过更强的运行框架基线",
        9.16,
        2.89,
        2.95,
        0.65,
        size=10,
        color="C7CFD8",
    )
    add_line(slide, 1.06, 4.46, 12.02, 4.46, color="516176", width=1.0)
    stages = [
        ("01", "开放资源"),
        ("02", "技能 Wiki"),
        ("03", "检索与组合"),
        ("04", "软件产物"),
    ]
    for idx, (n, label) in enumerate(stages):
        sx = 1.06 + idx * 2.92
        add_text(slide, n, sx, 4.70, 0.45, 0.27, size=8, color=CORAL, bold=True, font=FONT_EN)
        add_text(slide, label, sx, 5.03, 1.85, 0.36, size=13, color=WHITE, bold=True)
        if idx < len(stages) - 1:
            add_text(slide, "→", sx + 2.21, 4.96, 0.35, 0.35, size=18, color="77879B", font=FONT_EN)
    add_text(
        slide,
        "Yijia Fan et al.  ·  arXiv:2606.29538v4  ·  2026-07-17  ·  论文汇报",
        1.06,
        6.68,
        10.4,
        0.28,
        size=8.5,
        color="AAB5C2",
        font=FONT_EN,
    )
    add_notes(
        slide,
        """
        这篇论文讨论的是软件智能体如何获得可复用的程序性知识。作者不把技能理解为一段提示词，而是把它组织成包含文本、视觉、代码、元数据与来源信息的可执行条目。
        汇报将围绕三个问题展开：Resource2Skill 如何构建技能，技能访问是否真的提升产物质量，以及哪些组件带来主要收益。
        """,
    )

    # 2 — Bottleneck
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "多模态教程蕴含操作知识，却难以直接进入智能体记忆", 2, "Problem")
    add_text(slide, "人类资源保留了智能体最缺的程序性信号", 0.86, 1.30, 4.0, 0.35, size=15, color=NAVY, bold=True)
    resources = [
        ("视频", "操作时序、界面变化、视觉效果", CORAL_LIGHT, CORAL),
        ("代码", "工具调用、参数与可执行模式", BLUE_LIGHT, BLUE),
        ("文章", "概念、适用条件与解释", GOLD_LIGHT, GOLD),
        ("产物", "风格、布局与完成质量", TEAL_LIGHT, TEAL),
    ]
    for idx, (name, desc, fill, color) in enumerate(resources):
        yy = 1.86 + idx * 0.86
        add_rect(slide, 0.86, yy, 3.82, 0.65, fill=fill, line=fill, radius=True)
        add_text(slide, name, 1.08, yy + 0.11, 0.62, 0.30, size=12, color=color, bold=True, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, 1.80, yy + 0.10, 2.55, 0.34, size=10.5, color=INK, valign=MSO_ANCHOR.MIDDLE)
    add_chevron(slide, 5.05, 2.90, color=MUTED)
    add_rect(slide, 5.56, 1.72, 2.65, 1.48, fill=CORAL_LIGHT, line=CORAL_LIGHT, radius=True)
    add_text(slide, "直接放入上下文", 5.86, 1.95, 2.05, 0.32, size=13, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "昂贵 · 冗余 · 难扩展", 5.86, 2.42, 2.05, 0.30, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_rect(slide, 5.56, 3.58, 2.65, 1.48, fill=GOLD_LIGHT, line=GOLD_LIGHT, radius=True)
    add_text(slide, "压缩为纯文本摘要", 5.77, 3.81, 2.23, 0.32, size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "丢失时序 · 布局 · 交互", 5.76, 4.28, 2.27, 0.30, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_chevron(slide, 8.58, 2.90, color=MUTED)
    add_rect(slide, 9.05, 1.72, 3.24, 3.34, fill=WHITE, line=TEAL, radius=True, line_width=1.6)
    add_text(slide, "需要的中间表示", 9.37, 2.02, 2.57, 0.36, size=12, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "可检索\n可组合\n可执行\n可追踪", 9.59, 2.58, 2.15, 1.55, size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER, line_spacing=1.18)
    add_rect(slide, 0.86, 5.66, 11.43, 0.79, fill=NAVY, line=NAVY, radius=True)
    add_text(
        slide,
        "核心缺口：把高维、多源的人类经验转换成可复用的程序性记忆。",
        1.23,
        5.83,
        10.65,
        0.36,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_source(slide, "依据：paper.md S002–S004")
    add_notes(
        slide,
        """
        软件智能体需要操作软件、调用工具、检查中间结果并恢复失败。决定成败的往往不是事实知识，而是如何拆解目标、选择工具、检查状态以及修复错误。
        视频保留操作时序和视觉变化，代码保留可执行模式，文章和参考产物提供概念与风格。但把原始视频直接放入记忆成本很高，压缩成摘要又会丢失关键的动态与空间信息。
        """,
    )

    # 3 — Research question and contributions
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "Resource2Skill把“资源消费”转化为可维护的程序性记忆", 3, "Research question")
    add_text(
        slide,
        "能否把人类创作的多模态资源自动蒸馏为技能，\n并构建可扩展的软件智能体技能库？",
        0.95,
        1.30,
        6.0,
        1.18,
        size=20,
        color=NAVY,
        bold=True,
        line_spacing=1.02,
    )
    add_line(slide, 0.96, 2.72, 6.62, 2.72, color=CORAL, width=2.0)
    add_text(
        slide,
        "论文的主张覆盖构建与使用两个阶段，并用同一算子连接离线库与在线补缺。",
        0.96,
        3.08,
        5.72,
        0.82,
        size=13,
        color=SLATE,
    )
    stages = [
        ("01", "蒸馏", "从视频、代码、文章、参考产物提取技能"),
        ("02", "组织", "以分层、多模态 Skill Wiki 维护"),
        ("03", "选择", "MetaBrowse 收窄范围并选择技能子集"),
        ("04", "扩展", "覆盖不足时受控在线获取新技能"),
    ]
    for idx, (n, label, desc) in enumerate(stages):
        yy = 1.28 + idx * 1.22
        add_text(slide, n, 7.35, yy, 0.55, 0.30, size=9, color=CORAL, bold=True, font=FONT_EN)
        add_text(slide, label, 7.98, yy - 0.04, 1.00, 0.35, size=15, color=NAVY, bold=True)
        add_text(slide, desc, 9.12, yy - 0.01, 3.05, 0.52, size=11, color=INK)
        if idx < 3:
            add_line(slide, 7.62, yy + 0.46, 7.62, yy + 1.07, color=LINE, width=1.2)
    add_rect(slide, 0.96, 5.24, 5.73, 1.00, fill=TEAL_LIGHT, line=TEAL_LIGHT, radius=True)
    add_text(
        slide,
        "技能库并非固定提示词集合，而是一套可增长、可审计的长期记忆。",
        1.25,
        5.47,
        5.15,
        0.44,
        size=14,
        color="176C61",
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "依据：paper.md S003–S006")
    add_notes(
        slide,
        """
        研究问题是能否把多模态人类资源自动转换为可执行技能，并形成可扩展技能库。
        论文的贡献不是单一蒸馏模型，而是一条端到端闭环：构建条目、组织 Wiki、检索选择、在具体软件中执行，并在离线覆盖不足时使用同一构建算子在线补缺。
        """,
    )

    # 4 — Pipeline
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "系统将构建、组织、选择与执行连成同一条流水线", 4, "Method overview")
    add_rect(slide, 0.62, 1.18, 12.08, 5.62, fill=WHITE, line=LINE, radius=True)
    add_picture_contain(slide, ASSETS / "fig2_pipeline.png", 0.78, 1.31, 11.76, 5.25)
    add_source(slide, "Source: Fig. 2, Resource2Skill, arXiv:2606.29538v4 (2026)", y=6.86)
    add_notes(
        slide,
        """
        这张图是整篇论文的主干。左侧四类资源先经过解析与蒸馏，再通过完整性、去重、可执行性和安全性质量门控，进入 LM Wiki。
        技能条目包含文本、视觉、代码和元数据。推理时，智能体浏览 Wiki、选择技能并组合产物，领域适配器负责把技能转换为具体软件约束和动作。
        底部给出评测协议：相同测试任务下比较完整 Wiki、扁平技能库与无技能智能体，并对渲染产物评分。
        """,
    )

    # 5 — Skill representation
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "一项技能同时保留文本、视觉、代码、元数据与分类路径", 5, "Skill schema")
    add_text(slide, "s = ( p, x_text, x_visual, x_code, m )", 0.95, 1.16, 11.4, 0.53, size=23, color=NAVY, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.88, 2.06, 2.03, 3.66, fill=SOFT, line=LINE, radius=True)
    add_text(slide, "分类路径 p", 1.15, 2.36, 1.48, 0.36, size=14, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    for idx, label in enumerate(["领域", "类别", "子类别", "技能"]):
        yy = 3.00 + idx * 0.56
        add_text(slide, label, 1.24 + idx * 0.10, yy, 1.15, 0.28, size=11, color=INK, align=PP_ALIGN.CENTER)
        if idx < 3:
            add_text(slide, "↓", 1.68 + idx * 0.10, yy + 0.31, 0.25, 0.20, size=11, color=MUTED, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_chevron(slide, 3.13, 3.66, color=MUTED)
    add_rect(slide, 3.57, 1.90, 5.52, 4.02, fill=WHITE, line=TEAL, radius=True, line_width=1.6)
    add_text(slide, "Skill Wiki entry", 3.88, 2.14, 4.91, 0.38, size=14, color=TEAL, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)
    bands = [
        ("TEXT", "适用性、机制、输入与预期效果", BLUE_LIGHT, BLUE),
        ("VISUAL", "截图、布局、动效与渲染示例", CORAL_LIGHT, CORAL),
        ("CODE", "可执行或可改写的过程片段", GOLD_LIGHT, GOLD),
    ]
    for idx, (label, desc, fill, color) in enumerate(bands):
        yy = 2.76 + idx * 0.86
        add_rect(slide, 3.98, yy, 4.71, 0.65, fill=fill, line=fill, radius=True)
        add_text(slide, label, 4.22, yy + 0.13, 0.82, 0.30, size=10, color=color, bold=True, font=FONT_EN, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, desc, 5.26, yy + 0.12, 3.05, 0.32, size=10.5, color=INK, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "Linked · versioned · traceable", 4.12, 5.45, 4.44, 0.25, size=9.5, color=TEAL, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_chevron(slide, 9.36, 3.66, color=MUTED)
    add_rect(slide, 9.81, 2.07, 2.60, 1.38, fill=PURPLE_LIGHT, line=PURPLE_LIGHT, radius=True)
    add_text(slide, "元数据 m", 10.13, 2.31, 1.95, 0.32, size=14, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "标签 · 来源 · 层级 · 审计", 10.07, 2.80, 2.08, 0.31, size=10, color=SLATE, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.81, 3.77, 2.60, 1.77, fill=TEAL_LIGHT, line=TEAL_LIGHT, radius=True)
    add_text(slide, "质量门控 A_D", 10.08, 4.02, 2.04, 0.34, size=13.5, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "完整性 · 来源\n去重 · 模态一致\n代码结构可执行", 10.15, 4.50, 1.90, 0.76, size=10, color=INK, align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_rect(slide, 0.89, 6.16, 11.52, 0.48, fill=NAVY, line=NAVY, radius=True)
    add_text(slide, "同一条目同时回答：何时适用、如何执行、预期产生什么效果，以及证据来自哪里。", 1.12, 6.27, 11.02, 0.25, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "依据：paper.md S009–S010", y=6.86)
    add_notes(
        slide,
        """
        技能被表示为五元组。分类路径 p 让系统可以沿领域层级导航；元数据 m 支持过滤、审计和来源追踪。
        三种内容视图互补：文本说明何时使用，视觉保留布局与效果，代码给出工具化执行模式。候选技能还要通过完整性、来源、去重、模态一致与代码结构可执行性检查。
        """,
    )

    # 6 — MetaBrowse
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "MetaBrowse先收窄分类树，再由语言模型选择可组合技能", 6, "Inference")
    process = [
        ("任务 q", "用户需求"),
        ("词法收窄", "名称 · 标签 · 适用性 · 路径"),
        ("LM 选择", "读取文本 / 视觉 / 代码证据"),
        ("子集组合", "允许零项或多项技能"),
        ("MCP 执行", "应用 · 渲染 · 检查"),
    ]
    widths = [1.45, 2.15, 2.25, 1.85, 1.85]
    start_x = 0.77
    xx = start_x
    for idx, ((title, desc), ww) in enumerate(zip(process, widths)):
        fill = TEAL_LIGHT if idx in (2, 3) else WHITE
        line_color = TEAL if idx in (2, 3) else LINE
        add_rect(slide, xx, 1.69, ww, 1.47, fill=fill, line=line_color, radius=True, line_width=1.2)
        add_text(slide, f"0{idx + 1}", xx + 0.16, 1.86, 0.35, 0.22, size=8, color=CORAL, bold=True, font=FONT_EN)
        add_text(slide, title, xx + 0.16, 2.16, ww - 0.32, 0.34, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, xx + 0.14, 2.62, ww - 0.28, 0.34, size=8.8, color=SLATE, align=PP_ALIGN.CENTER)
        xx += ww
        if idx < len(process) - 1:
            add_chevron(slide, xx + 0.04, 2.18, color=MUTED)
            xx += 0.43
    add_line(slide, 1.39, 3.65, 11.96, 3.65, color=LINE, width=0.9)
    add_text(slide, "若候选集合不足", 1.03, 4.11, 2.00, 0.33, size=12, color=CORAL, bold=True)
    add_text(slide, "↓", 1.75, 4.46, 0.25, 0.27, size=15, color=CORAL, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.91, 4.78, 2.15, 1.01, fill=CORAL_LIGHT, line=CORAL_LIGHT, radius=True)
    add_text(slide, "在线构建算子", 1.19, 4.98, 1.59, 0.32, size=13, color=CORAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "搜索 · 蒸馏 · 验收", 1.20, 5.39, 1.59, 0.25, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_chevron(slide, 3.31, 5.04, color=MUTED)
    add_rect(slide, 3.76, 4.78, 2.37, 1.01, fill=PURPLE_LIGHT, line=PURPLE_LIGHT, radius=True)
    add_text(slide, "独立在线技能池", 4.00, 4.98, 1.89, 0.32, size=13, color=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "不回写离线 Wiki", 4.07, 5.39, 1.74, 0.25, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_chevron(slide, 6.37, 5.04, color=MUTED)
    add_rect(slide, 6.83, 4.78, 2.60, 1.01, fill=TEAL_LIGHT, line=TEAL_LIGHT, radius=True)
    add_text(slide, "重新选择与组合", 7.13, 4.98, 2.00, 0.32, size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "沿用同一执行接口", 7.22, 5.39, 1.82, 0.25, size=9.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "子集选择 ≠ 排名", 10.22, 4.79, 1.89, 0.32, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "LM 可以选零项技能，\n也可组合互补技能。", 10.24, 5.29, 1.86, 0.67, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_source(slide, "依据：paper.md S011–S012、S024")
    add_notes(
        slide,
        """
        MetaBrowse 不是把全部技能做一次相似度排序。第一阶段用名称、标签、适用性文本和分类路径定位相关子树；第二阶段由语言模型读取候选条目的结构化证据并选择一个子集。
        这里选择的是子集而不是单一排名，因此可以不选、选一项或组合多项技能。若离线候选不足，系统在线搜索并蒸馏新技能，放入独立在线池后再执行。
        """,
    )

    # 7 — Evaluation design
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "评测用配对任务隔离技能访问与执行接口的贡献", 7, "Evaluation")
    flow = [
        ("配对任务", "每领域 N=80\n消融 N=40", 1.78),
        ("4 种后端", "GPT-5.5 / 5.4\nMini / Nano", 2.00),
        ("4 种系统", "w Skills / w/o Skills\nClaudeCode-H / Codex-H", 2.50),
        ("渲染产物", "同任务、裁判与随机种子", 2.10),
        ("盲测裁判", "视觉模型；Reaper 为音频模型", 2.15),
    ]
    xx = 0.70
    for idx, (title, desc, ww) in enumerate(flow):
        fill = TEAL_LIGHT if idx == 2 else WHITE
        line_color = TEAL if idx == 2 else LINE
        add_rect(slide, xx, 1.44, ww, 1.55, fill=fill, line=line_color, radius=True, line_width=1.2)
        add_text(slide, title, xx + 0.16, 1.69, ww - 0.32, 0.34, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, xx + 0.14, 2.18, ww - 0.28, 0.58, size=9.4, color=SLATE, align=PP_ALIGN.CENTER)
        xx += ww
        if idx < len(flow) - 1:
            add_chevron(slide, xx + 0.03, 1.98, color=MUTED)
            xx += 0.42
    add_text(slide, "覆盖 7 个软件创作领域", 0.84, 3.53, 2.63, 0.36, size=14, color=NAVY, bold=True)
    domains = ["Web", "Excel", "Reaper", "PPT", "Blender", "CAD", "UE5"]
    colors = [BLUE, TEAL, PURPLE, CORAL, GOLD, SLATE, NAVY]
    for idx, (domain, color) in enumerate(zip(domains, colors)):
        xx = 0.86 + idx * 1.72
        add_line(slide, xx, 4.15, xx + 1.18, 4.15, color=color, width=3.2)
        add_text(slide, domain, xx, 4.32, 1.18, 0.31, size=10.8, color=color, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.86, 5.00, 5.85, 1.12, fill=SOFT, line=SOFT, radius=True)
    add_text(slide, "主指标", 1.13, 5.22, 0.80, 0.29, size=11.5, color=CORAL, bold=True)
    add_text(slide, "各领域五维量表的无权平均；未产出可评分产物记 0 分", 2.08, 5.18, 4.30, 0.52, size=11, color=INK)
    add_rect(slide, 7.02, 5.00, 5.18, 1.12, fill=BLUE_LIGHT, line=BLUE_LIGHT, radius=True)
    add_text(slide, "人类验证", 7.29, 5.22, 1.00, 0.29, size=11.5, color=BLUE, bold=True)
    add_text(slide, "40 组匹配样本 × 每组 5 位评分者，匿名并排 A/B", 8.43, 5.18, 3.45, 0.52, size=11, color=INK)
    add_source(slide, "依据：paper.md S013–S015、S026")
    add_notes(
        slide,
        """
        评测覆盖七个创作领域。主比较在每个领域使用匹配的 N=80 任务，并在四种后端和四种系统配置上重复；消融使用 N=40。
        同一个模型×领域单元内固定任务、裁判和随机种子，目的是把差异归因于执行接口与技能访问。
        主要分数来自模型裁判，同时作者用 40 组样本、每组 5 位评分者的人类盲测验证偏好方向。
        """,
    )

    # 8 — Main result
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "技能访问在全部 28 个“模型×领域”单元中带来提升", 8, "Main result")
    add_rect(slide, 0.61, 1.18, 8.78, 5.66, fill=WHITE, line=LINE, radius=True)
    add_picture_contain(slide, ASSETS / "table1_main_comparison.png", 0.72, 1.30, 8.55, 5.43)
    add_metric(slide, "+11.9 pp", "平均：56.8% vs 45.0%", 9.75, 1.38, 2.66, color=CORAL, value_size=28)
    add_line(slide, 9.76, 2.48, 12.14, 2.48, color=LINE, width=1.0)
    add_metric(slide, "28 / 28", "w Skills 均优于 w/o Skills", 9.75, 2.71, 2.66, color=TEAL, value_size=25)
    add_line(slide, 9.76, 3.80, 12.14, 3.80, color=LINE, width=1.0)
    add_metric(slide, "26 / 28", "超过更强的两种 harness", 9.75, 4.03, 2.66, color=BLUE, value_size=25)
    add_line(slide, 9.76, 5.11, 12.14, 5.11, color=LINE, width=1.0)
    add_metric(slide, "p < 10⁻³", "论文报告的配对 Wilcoxon 检验", 9.75, 5.35, 2.66, color=PURPLE, value_size=19)
    add_source(slide, "Source: Table 1; paper.md S015–S016", y=6.87)
    add_notes(
        slide,
        """
        主比较显示，w Skills 在 28 个模型×领域单元中全部超过同一智能体的无技能版本，整体平均从 45.0% 提高到 56.8%，即 11.9 个百分点。
        两种现成运行框架也能改善无技能智能体，但完整 Resource2Skill 在 28 个单元中的 26 个仍超过更强的基线，两个例外差距不到 1 分。
        作者报告配对 Wilcoxon 检验达到 p 小于 10 的负三次方。
        """,
    )

    # 9 — Domain and human evidence
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "收益集中于复杂创作软件，人类盲测支持同一偏好方向", 9, "Where gains appear")
    add_rect(slide, 0.60, 1.22, 9.35, 4.98, fill=WHITE, line=LINE, radius=True)
    add_picture_contain(slide, ASSETS / "fig1_overview.png", 0.72, 1.34, 9.10, 4.72)
    add_metric(slide, "+30–40 pp", "UE5 的技能增益最大", 10.22, 1.43, 2.20, color=CORAL, value_size=24)
    add_line(slide, 10.22, 2.56, 12.27, 2.56, color=LINE, width=1.0)
    add_metric(slide, "85.5%", "排除平局的人类偏好胜率", 10.22, 2.82, 2.20, color=TEAL, value_size=25)
    add_line(slide, 10.22, 3.95, 12.27, 3.95, color=LINE, width=1.0)
    add_metric(slide, "α = 0.58", "评分者一致性为中等水平", 10.22, 4.22, 2.20, color=PURPLE, value_size=19)
    add_rect(slide, 0.71, 6.35, 11.66, 0.42, fill=NAVY, line=NAVY, radius=True)
    add_text(slide, "工具接口与操作规范越密集，技能越能减少从提示词重新推导流程的成本。", 0.96, 6.43, 11.17, 0.24, size=11.7, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: Fig. 1; paper.md S017、S026", y=6.87)
    add_notes(
        slide,
        """
        分领域看，Blender、Web 和 UE5 的提升最大，尤其 UE5 达到 30 到 40 个百分点。论文解释为这些工具的创作规范和接口细节更密集，难以从单次提示重新推导。
        人类盲测中，w Skills 在 200 次个体评分里获胜 136 次、平局 41 次；排除平局后胜率为 85.5%。但 Krippendorff α 为 0.58，说明评分者一致性只有中等水平。
        """,
    )

    # 10 — Scaling and wiki
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "前 200 项技能贡献最大，完整 Wiki 继续带来结构性增益", 10, "Scaling")
    add_rect(slide, 0.62, 1.19, 12.08, 4.62, fill=WHITE, line=LINE, radius=True)
    add_picture_contain(slide, ASSETS / "fig3_scaling_wiki.png", 0.74, 1.34, 11.84, 4.30)
    facts = [
        ("0 → 200", "各领域最大增益区间"),
        ("400 → Full", "单领域最多再增加 0.8 pp"),
        ("Our Wiki", "持续优于扁平纯文本访问"),
    ]
    for idx, (value, label) in enumerate(facts):
        xx = 0.90 + idx * 4.08
        add_text(slide, value, xx, 6.04, 1.54, 0.34, size=15, color=[CORAL, GOLD, TEAL][idx], bold=True, font=FONT_EN)
        add_text(slide, label, xx + 1.60, 6.07, 2.23, 0.31, size=9.8, color=SLATE)
    add_source(slide, "Source: Fig. 3; paper.md S018", y=6.87)
    add_notes(
        slide,
        """
        随技能库扩展，各领域分数单调上升，并在约 200 项技能后趋于饱和。0 到 200 项贡献最大，Excel 增加 14.2 个百分点，Reaper 增加 3.1 个百分点。
        从 400 项到完整库，每个领域至多增加 0.8 分，说明后续条目主要补充长尾能力。
        右侧实验显示，纯文本技能访问已经有效，但分层 Wiki、视觉和代码信息继续提升表现。
        """,
    )

    # 11 — Online acquisition
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "在线获取只在离线技能池存在覆盖缺口时显著奏效", 11, "Online acquisition")
    add_text(slide, "离线池固定为 891 项技能；在线分支最多新增 100 项", 0.94, 1.24, 6.4, 0.36, size=13, color=SLATE)
    rows = [
        ("T_standard", "常规任务", 65.4, 66.1, "+0.7 pp", BLUE),
        ("T_novel", "预先确认的覆盖缺口", 41.2, 62.8, "+21.6 pp", CORAL),
    ]
    for idx, (task, label, before, after, delta, color) in enumerate(rows):
        yy = 1.86 + idx * 2.05
        add_text(slide, task, 0.96, yy, 1.48, 0.36, size=15, color=NAVY, bold=True, font=FONT_EN)
        add_text(slide, label, 0.96, yy + 0.42, 2.06, 0.32, size=10.5, color=SLATE)
        add_text(slide, f"{before:.1f}", 3.20, yy + 0.15, 0.82, 0.42, size=22, color=MUTED, bold=True, font=FONT_EN, align=PP_ALIGN.RIGHT)
        add_text(slide, "offline", 3.25, yy + 0.63, 0.72, 0.22, size=8, color=MUTED, font=FONT_EN, align=PP_ALIGN.RIGHT)
        add_line(slide, 4.35, yy + 0.39, 8.56, yy + 0.39, color=LINE, width=6.0)
        frac_before = before / 70.0
        frac_after = after / 70.0
        add_line(slide, 4.35, yy + 0.39, 4.35 + 4.21 * frac_before, yy + 0.39, color="AAB4BE", width=6.0)
        add_line(slide, 4.35, yy + 0.79, 4.35 + 4.21 * frac_after, yy + 0.79, color=color, width=6.0)
        add_text(slide, f"{after:.1f}", 8.83, yy + 0.19, 0.78, 0.42, size=22, color=color, bold=True, font=FONT_EN)
        add_text(slide, "offline + online", 8.84, yy + 0.65, 1.20, 0.22, size=8, color=MUTED, font=FONT_EN)
        add_rect(slide, 10.52, yy + 0.09, 1.65, 0.83, fill=CORAL_LIGHT if idx else BLUE_LIGHT, line=CORAL_LIGHT if idx else BLUE_LIGHT, radius=True)
        add_text(slide, delta, 10.68, yy + 0.28, 1.33, 0.35, size=18, color=color, bold=True, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.94, 6.05, 11.27, 0.62, fill=NAVY, line=NAVY, radius=True)
    add_text(slide, "作者将在线获取定位为受控补缺，并在常规基准中默认关闭。", 1.22, 6.19, 10.70, 0.31, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: Table 2; paper.md S019、S024", y=6.87)
    add_notes(
        slide,
        """
        在线获取实验区分常规任务和离线池明确缺失能力的任务。常规任务中新增 100 项技能只提高 0.7 分；在 T_novel 中，分数从 41.2% 提高到 62.8%，增加 21.6 分。
        因此，这个结果支持在线技能作为覆盖缺口补充，而不能推断在线搜索对自然任务分布总有正收益。在线获取还增加搜索、蒸馏和验证延迟。
        """,
    )

    # 12 — Source ablation
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "视频提供不可替代的时序与视觉操作信号", 12, "Source ablation")
    add_rect(slide, 0.61, 1.25, 8.71, 4.35, fill=WHITE, line=LINE, radius=True)
    add_picture_contain(slide, ASSETS / "table3_source_ablation.png", 0.74, 1.38, 8.45, 4.08)
    add_metric(slide, "68.9%", "完整四来源技能池", 9.68, 1.44, 2.55, color=TEAL, value_size=25)
    add_line(slide, 9.69, 2.53, 12.03, 2.53, color=LINE, width=1.0)
    add_metric(slide, "59.4%", "移除视频：下降 9.5 pp", 9.68, 2.77, 2.55, color=CORAL, value_size=25)
    add_line(slide, 9.69, 3.86, 12.03, 3.86, color=LINE, width=1.0)
    add_metric(slide, "66.8%", "仅视频仍高于无视频三来源 7.4 分", 9.68, 4.09, 2.55, color=GOLD, value_size=25)
    add_rect(slide, 0.87, 5.92, 4.82, 0.55, fill=CORAL_LIGHT, line=CORAL_LIGHT, radius=True)
    add_text(slide, "Excel  −14.2 pp", 1.12, 6.04, 1.79, 0.28, size=13, color=CORAL, bold=True, font=FONT_EN)
    add_text(slide, "Web  −11.5 pp", 3.53, 6.04, 1.66, 0.28, size=13, color=CORAL, bold=True, font=FONT_EN)
    add_text(slide, "去除视频的损失集中在依赖操作顺序和视觉变化的领域。", 6.05, 5.98, 6.05, 0.48, size=12, color=SLATE)
    add_source(slide, "Source: Table 3; paper.md S020", y=6.87)
    add_notes(
        slide,
        """
        来源消融中，移除视频后平均分从 68.9% 降到 59.4%。仅使用视频的技能库达到 66.8%，仍比代码、文章和参考产物三来源但无视频的配置高 7.4 分。
        去除视频在 Excel 和 Web 上损失最大，分别下降 14.2 和 11.5 个百分点，说明操作时序与视觉变化难以被纯文本和代码完全替代。
        """,
    )

    # 13 — Representation ablation
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "整理后的文本已贡献大部分收益，多模态内容再增加 3.9 分", 13, "Representation ablation")
    data = [
        ("Text", 65.0),
        ("Text + Visual", 66.9),
        ("Text + Code", 67.0),
        ("Full", 68.9),
    ]
    horizontal_bars(slide, data, 0.88, 1.50, 7.15, 3.78, 60.0, 70.0, highlight_index=3, label_w=1.80)
    add_text(slide, "横轴：总体得分 60–70%（截断坐标）", 2.66, 5.47, 3.66, 0.25, size=8.2, color=MUTED, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_line(slide, 8.50, 1.44, 8.50, 5.88, color=LINE, width=1.0)
    add_text(slide, "如何理解这 3.9 分", 8.91, 1.53, 2.82, 0.38, size=16, color=NAVY, bold=True)
    add_bullets(
        slide,
        [
            "Text 已包含适用性与路由线索",
            "Visual 在 Text 上增加 1.9 pp",
            "Code 在 Text 上增加 2.0 pp",
            "Full 在每个领域均排名第一",
        ],
        8.90,
        2.15,
        3.25,
        2.38,
        size=12,
        color=INK,
        gap=10,
    )
    add_rect(slide, 8.89, 4.89, 3.11, 0.78, fill=TEAL_LIGHT, line=TEAL_LIGHT, radius=True)
    add_text(slide, "整体 11.9 分并非全部来自多模态表示；\n表 4 隔离出的独立增益约为 3.9 分。", 9.13, 5.03, 2.65, 0.49, size=10.5, color="176C61", bold=True, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: Table 4; paper.md S021", y=6.87)
    add_notes(
        slide,
        """
        匹配预算的表示消融显示，文本配置已经达到 65.0%，因为它包含人工或模型整理后的适用性和路由信息。
        在此基础上加入视觉提高 1.9 分，加入代码提高 2.0 分，完整多模态条目达到 68.9%。
        因此需要区分两件事：端到端技能系统带来的 11.9 分，与多模态表示在整理文本之上的约 3.9 分独立增益。
        """,
    )

    # 14 — Selection ablation
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "MetaBrowse胜过纯检索，说明“选什么、怎么组合”同样关键", 14, "Selection ablation")
    selection = [
        ("MetaBrowse", 68.9),
        ("BM25", 66.0),
        ("BM25 + Embed", 64.2),
        ("Embed", 60.0),
        ("Random-FullPool", 58.0),
        ("No-Skill", 57.3),
    ]
    horizontal_bars(slide, selection, 0.79, 1.35, 8.20, 4.70, 55.0, 70.0, highlight_index=0, label_w=2.05)
    add_text(slide, "横轴：总体得分 55–70%（截断坐标）", 3.00, 6.18, 3.80, 0.24, size=8.2, color=MUTED, font=FONT_EN, align=PP_ALIGN.CENTER)
    add_line(slide, 9.36, 1.45, 9.36, 5.96, color=LINE, width=1.0)
    add_text(slide, "相对最强纯检索基线", 9.76, 1.50, 2.32, 0.35, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_metric(slide, "+5.0 pp", "Excel", 9.83, 2.05, 2.20, color=CORAL, value_size=21)
    add_metric(slide, "+3.8 pp", "PPT", 9.83, 3.22, 2.20, color=TEAL, value_size=21)
    add_metric(slide, "+2.3 pp", "Blender", 9.83, 4.39, 2.20, color=PURPLE, value_size=21)
    add_text(slide, "BM25+Embed 未超过 BM25，\n向量相似度没有自动解决技能互补性。", 9.70, 5.60, 2.49, 0.61, size=10.5, color=SLATE, align=PP_ALIGN.CENTER)
    add_source(slide, "Source: Table 5; paper.md S022", y=6.87)
    add_notes(
        slide,
        """
        选择策略消融中，先沿分类层级收窄范围、再由语言模型选择子集的 MetaBrowse 达到 68.9%，高于 BM25 的 66.0%。
        BM25 加向量检索并没有超过纯 BM25，说明相似度检索不能自动解决技能是否适配、不同技能是否互补的问题。
        MetaBrowse 相对最强纯检索基线在 Excel、PPT 和 Blender 上的优势最大。
        """,
    )

    # 15 — Synthesis and limitations
    slide = prs.slides.add_slide(blank)
    add_standard_title(slide, "结构化技能记忆的价值已被支持，外推仍受三类边界约束", 15, "Synthesis")
    add_text(slide, "证据支持的结论", 0.88, 1.31, 4.80, 0.38, size=17, color=TEAL, bold=True)
    conclusions = [
        ("+11.9 pp", "完整技能流水线稳定改善七类软件产物"),
        ("复杂领域", "UE5、Blender、Web 的收益更集中"),
        ("+21.6 pp", "在线技能对已知覆盖缺口有效"),
    ]
    for idx, (value, desc) in enumerate(conclusions):
        yy = 1.95 + idx * 1.10
        add_text(slide, value, 0.93, yy, 1.46, 0.37, size=18, color=[CORAL, BLUE, PURPLE][idx], bold=True, font=FONT_EN if "+" in value else FONT_CN)
        add_text(slide, desc, 2.55, yy + 0.02, 3.42, 0.45, size=12, color=INK, bold=True)
        add_line(slide, 0.93, yy + 0.68, 5.95, yy + 0.68, color=LINE, width=0.8)
    add_line(slide, 6.38, 1.32, 6.38, 5.66, color=LINE, width=1.0)
    add_text(slide, "需要继续检验的边界", 6.82, 1.31, 4.85, 0.38, size=17, color=CORAL, bold=True)
    limits = [
        ("测量", "主要依赖模型裁判；人类一致性 α=0.58"),
        ("对照", "缺少等 token 预算的原始资源直接检索基线"),
        ("外推", "依赖可编程工具接口、公开程序性内容，并承担在线延迟"),
    ]
    for idx, (label, desc) in enumerate(limits):
        yy = 1.93 + idx * 1.18
        add_text(slide, f"0{idx + 1}", 6.85, yy, 0.40, 0.27, size=8.5, color=CORAL, bold=True, font=FONT_EN)
        add_text(slide, label, 7.42, yy - 0.03, 0.66, 0.32, size=13, color=NAVY, bold=True)
        add_text(slide, desc, 8.23, yy - 0.02, 3.81, 0.59, size=11, color=INK)
    add_rect(slide, 0.86, 5.86, 11.45, 0.72, fill=NAVY, line=NAVY, radius=True)
    add_text(
        slide,
        "可复用的核心对象，是带来源、质量门控与执行接口的程序性记忆。",
        1.19,
        6.02,
        10.78,
        0.36,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_source(slide, "依据：paper.md S023、S025–S027 与批判性阅读提示", y=6.87)
    add_notes(
        slide,
        """
        最稳妥的结论是：结构化技能库可以改善复杂软件创作，而且在工具接口和操作规范密集的领域收益更大。
        但 11.9 分是完整系统效果，不能全部归因于多模态；多模态表示的匹配预算消融增益约 3.9 分。
        外推需要注意三点：分数主要来自模型裁判；论文没有与等预算的原始资源直接检索比较；方法依赖可编程工具接口和公开程序性资源，在线补缺还增加延迟。
        """,
    )

    prs.core_properties.title = "Resource2Skill：从人类创作资源中蒸馏智能体可执行技能"
    prs.core_properties.subject = "中文论文汇报"
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.comments = "Generated from resource2skill_reader/paper.md"
    prs.save(PPTX_PATH)
    return PPTX_PATH


if __name__ == "__main__":
    result = build_deck()
    print(result)
